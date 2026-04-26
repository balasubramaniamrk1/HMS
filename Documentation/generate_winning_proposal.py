import os
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- Configuration ---
PROJECT_DIR = os.path.dirname(os.path.abspath(__file__))
SCREENSHOT_DIR = os.path.join(PROJECT_DIR, "screenshots")
OUTPUT_PPTX = os.path.join(PROJECT_DIR, "HMS_Client_Proposal_Winning.pptx")

# Brand Colors (LUX Theme)
COLOR_PRIMARY = RGBColor(10, 37, 64)    # Deep Navy Blue
COLOR_ACCENT = RGBColor(212, 175, 55)   # Gold
COLOR_TEXT_STD = RGBColor(51, 51, 51)   # Dark Gray
COLOR_WHITE = RGBColor(255, 255, 255)

def apply_text_style(run, font_name='Segoe UI', font_size=18, color=COLOR_TEXT_STD, bold=False):
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # Background shape for a premium feel
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.color.rgb = COLOR_PRIMARY

    # Add a gold accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.5), prs.slide_width, Inches(0.1)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_ACCENT
    accent.line.color.rgb = COLOR_ACCENT

    # Add Custom Title Text Box
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2))
    tf = txBox.text_frame
    tf.clear()  # Clear default paragraph
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "NEXT-GENERATION HOSPITAL\nMANAGEMENT SYSTEM"
    apply_text_style(run, font_name='Segoe UI', font_size=48, color=COLOR_WHITE, bold=True)

    # Subtitle
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = f"\nComprehensive Enterprise Proposal\n{datetime.date.today().strftime('%B %d, %Y')}"
    apply_text_style(run2, font_name='Segoe UI', font_size=28, color=COLOR_ACCENT, bold=False)

def add_content_slide(prs, title_text, bullet_points):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title_text
    title_shape.text_frame.paragraphs[0].font.name = 'Segoe UI'
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        p.level = 0
        p.font.name = 'Segoe UI'
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_TEXT_STD

def add_screenshot_slide(prs, title_text, img_filename, description):
    # Blank layout is usually index 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title Bar
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = COLOR_PRIMARY
    title_bg.line.color.rgb = COLOR_PRIMARY
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), prs.slide_width - Inches(1), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    apply_text_style(run, font_name='Segoe UI', font_size=36, color=COLOR_WHITE, bold=True)

    # Description Box (bottom)
    descBox = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(1), prs.slide_width - Inches(1), Inches(0.8))
    desc_tf = descBox.text_frame
    desc_p = desc_tf.paragraphs[0]
    desc_p.alignment = PP_ALIGN.CENTER
    desc_run = desc_p.add_run()
    desc_run.text = description
    apply_text_style(desc_run, font_name='Segoe UI', font_size=20, color=COLOR_TEXT_STD, bold=True)
    
    # Image logic for 16:9
    img_path = os.path.join(SCREENSHOT_DIR, img_filename)
    if os.path.exists(img_path):
        max_height = Inches(6.8)
        
        # Insert image fitting the height
        pic = slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.1), height=max_height)
        
        # Center horizontally
        if pic.width < prs.slide_width:
             pic.left = int((prs.slide_width - pic.width) / 2)
    else:
        print(f"Warning: Missing screenshot {img_filename}")
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(14), Inches(5))
        box.text = f"[ Missing Image: {img_filename} ]"

def create_presentation():
    # Standard aspect ratio 16:9 for enterprise.
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    print("Generating comprehensive enterprise presentation...")

    # 1. Main Title
    add_title_slide(prs)

    # 2. Executive Overview
    add_content_slide(prs, "The Legacy Challenge", [
        "Disconnected systems leading to operational silos across departments.",
        "High margin of error in patient records, prescriptions, and billing.",
        "Lack of real-time monitoring of hospital resources, beds, and staff.",
        "Poor patient experience due to fragmented service delivery and long wait times."
    ])

    add_content_slide(prs, "Our Enterprise Solution", [
        "A unified, modern Hospital Management System (HMS).",
        "Seamless digital integration from Admissions to Discharge and Billing.",
        "Real-time, responsive web interface tailored for doctors, nurses, and admins.",
        "Role-based access control ensuring patient data privacy (HIPAA compliant architecture).",
        "Comprehensive dashboard and analytics for hospital executives."
    ])

    # 3. Core Module Showcases
    add_screenshot_slide(prs, "Executive Dashboard", "01_hms_dashboard.png", 
        "A high-level overview of total admissions, active appointments, and hospital revenue metrics.")
    
    add_screenshot_slide(prs, "Admissions Management", "02_admissions_dashboard.png", 
        "Streamlined patient admission, discharge, and transfer (ADT) workflows.")
    
    add_screenshot_slide(prs, "Admit Patient Flow", "03_admissions_admit.png", 
        "Rapid patient intake forms ensuring accurate demographic and medical capture.")

    add_screenshot_slide(prs, "Appointments Dashboard", "04_appointments_staff.png", 
        "Robust appointment scheduling tools for Frontdesk and Nursing staff.")
    
    add_screenshot_slide(prs, "Doctor Console", "05_appointments_doctor.png", 
        "Specialized, distraction-free console for Doctors to manage consultations and prescriptions.")
    
    add_screenshot_slide(prs, "Pharmacy POS", "08_pharmacy_pos.png", 
        "Lightning-fast Point-of-Sale interface designed for high-volume outpatient pharmacies.")

    add_screenshot_slide(prs, "Pharmacy Purchase Orders", "09_pharmacy_purchase_orders.png", 
        "Automated PO generation and tracking for critical medicine stock management.")
        
    add_screenshot_slide(prs, "Pharmacy Returns", "10_pharmacy_returns.png", 
        "Auditable tracking of vendor returns and expired medication disposal.")

    add_screenshot_slide(prs, "Enterprise Inventory", "11_inventory_dashboard.png", 
        "Dashboard tracking overall hospital assets, AMCs, and equipment health.")
    
    add_screenshot_slide(prs, "Asset Tracking", "12_inventory_list.png", 
        "Detailed registry of hospital equipment including QR code capabilities.")
        
    add_screenshot_slide(prs, "Vendor Management", "13_inventory_vendors.png", 
        "Centralized directory for all medical and equipment suppliers.")

    add_screenshot_slide(prs, "Doctor Directory", "14_doctors_list.png", 
        "Organized staff directory grouped by medical specialties.")
        
    add_screenshot_slide(prs, "Departments & Specialties", "15_departments.png", 
        "Management of hospital Centers of Excellence and their capabilities.")

    add_screenshot_slide(prs, "Staff Attendance", "16_staff_attendance.png", 
        "Clock-in/Clock-out system monitoring daily employee attendance.")
    
    add_screenshot_slide(prs, "Staff Administration", "17_staff_admin.png", 
        "HR portal for managing employee roles, details, and access levels.")

    add_screenshot_slide(prs, "System Administration", "18_admin_console.png", 
        "Centralized Django admin panel for complete system configuration and metadata management.")

    # 4. Next Steps
    add_content_slide(prs, "Implementation Strategy", [
        "Phase 1: Environment Setup & Data Migration (Weeks 1-2)",
        "Phase 2: Master Data Configuration & Core Module Rollout (Weeks 3-4)",
        "Phase 3: Staff Training & Parallel Run (Week 5)",
        "Phase 4: Full System Go-Live and Post-Launch Support (Week 6 onwards)"
    ])

    prs.save(OUTPUT_PPTX)
    print(f"Presentation saved successfully to {OUTPUT_PPTX}")

if __name__ == "__main__":
    create_presentation()
